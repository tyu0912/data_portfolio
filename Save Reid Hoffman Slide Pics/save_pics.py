import urllib.request
from bs4 import BeautifulSoup
import re
from pptx import Presentation
from pptx.util import Inches
from io import BytesIO
from PIL import Image
import requests


prs = Presentation()
title_slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(title_slide_layout)

url = 'http://www.businessinsider.com/career-advice-from-billionaire-hoffman-2015-1?op=1'
html = urllib.request.urlopen(url).read()
soup = BeautifulSoup(html, "html.parser")

img = soup.findAll('img', {'data-src': re.compile("^[http://static4.businessinsider.com/]")})

count = 1

for link in img:
    link = str(link)
    thebeg = link.find('data-src="')
    end = link.find('src="', thebeg+15)
    final = link[thebeg+len('data-src="'):end-2]
    
    finalname = "Reid Hoffman Slides/" + str(count) + '.jpg'
    
    count = count + 1 
    
    urllib.request.urlretrieve(final, finalname)

    
    #slideimg = io.BytesIO(urllib.request.urlopen(final).read())
    #theimg = Image.open(slideimg)
    
    # response = requests.get(final).text
    # #theimg = Image.open(BytesIO(response.content))
    # theimg = Image.open(response.content)
    
    # left = top = Inches(1)
    # pic = slide.shapes.add_picture(theimg, left, top)
    
    # slide = prs.slides.add_slide(blank_slide_layout)

# prs.save('test.pptx')


